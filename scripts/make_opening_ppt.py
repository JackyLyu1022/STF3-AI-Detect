from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

OUT = Path('docs/STF3_Detect_opening_defense.pptx')
ASSETS = Path('docs/ppt_assets')
prs = Presentation()
prs.slide_width = Inches(13.333333)
prs.slide_height = Inches(7.5)

NAVY=RGBColor(23,44,82); BLUE=RGBColor(41,98,255); CYAN=RGBColor(0,172,193); ORANGE=RGBColor(245,124,0); GREEN=RGBColor(46,125,50); GRAY=RGBColor(90,100,115); LIGHT=RGBColor(245,248,252); WHITE=RGBColor(255,255,255); RED=RGBColor(198,40,40)
FONT='Microsoft YaHei'

def bg(slide):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb=WHITE
    top=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,Inches(0.12)); top.fill.solid(); top.fill.fore_color.rgb=BLUE; top.line.fill.background()
    foot=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,Inches(7.22),prs.slide_width,Inches(0.28)); foot.fill.solid(); foot.fill.fore_color.rgb=LIGHT; foot.line.fill.background()

def title(slide, t, sub=None):
    box=slide.shapes.add_textbox(Inches(0.55),Inches(0.34),Inches(12.2),Inches(0.55)); p=box.text_frame.paragraphs[0]
    p.text=t; p.font.name=FONT; p.font.size=Pt(26); p.font.bold=True; p.font.color.rgb=NAVY
    if sub:
        b=slide.shapes.add_textbox(Inches(0.58),Inches(0.88),Inches(12),Inches(0.3)); p=b.text_frame.paragraphs[0]
        p.text=sub; p.font.name=FONT; p.font.size=Pt(12); p.font.color.rgb=GRAY

def page(slide,n):
    b=slide.shapes.add_textbox(Inches(12.15),Inches(7.23),Inches(0.8),Inches(0.22)); p=b.text_frame.paragraphs[0]
    p.text=f'{n:02d}'; p.font.name='Arial'; p.font.size=Pt(10); p.font.color.rgb=GRAY; p.alignment=PP_ALIGN.RIGHT

def bullets(slide, items, x,y,w,h, size=18, color=NAVY):
    b=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=b.text_frame; tf.clear(); tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text='• '+it; p.font.name=FONT; p.font.size=Pt(size); p.font.color.rgb=color; p.space_after=Pt(8)

def card(slide,x,y,w,h,t,body,color=BLUE):
    s=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h)); s.fill.solid(); s.fill.fore_color.rgb=LIGHT; s.line.color.rgb=RGBColor(220,228,240)
    bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(0.08),Inches(h)); bar.fill.solid(); bar.fill.fore_color.rgb=color; bar.line.fill.background()
    b=slide.shapes.add_textbox(Inches(x+0.18),Inches(y+0.12),Inches(w-0.3),Inches(0.35)); p=b.text_frame.paragraphs[0]; p.text=t; p.font.name=FONT; p.font.size=Pt(15); p.font.bold=True; p.font.color.rgb=color
    b=slide.shapes.add_textbox(Inches(x+0.18),Inches(y+0.53),Inches(w-0.35),Inches(h-0.6)); p=b.text_frame.paragraphs[0]; p.text=body; p.font.name=FONT; p.font.size=Pt(12); p.font.color.rgb=NAVY

def pic(slide,path,x,y,w=None,h=None):
    path=Path(path)
    if path.exists():
        if w and h: slide.shapes.add_picture(str(path),Inches(x),Inches(y),Inches(w),Inches(h))
        elif w: slide.shapes.add_picture(str(path),Inches(x),Inches(y),width=Inches(w))
        elif h: slide.shapes.add_picture(str(path),Inches(x),Inches(y),height=Inches(h))
    else: card(slide,x,y,w or 4,h or 2,'Image missing',str(path),RED)

def table(slide,data,x,y,w,h,size=12):
    tbl=slide.shapes.add_table(len(data),len(data[0]),Inches(x),Inches(y),Inches(w),Inches(h)).table
    for c in range(len(data[0])): tbl.columns[c].width=Inches(w/len(data[0]))
    for r,row in enumerate(data):
        for c,val in enumerate(row):
            cell=tbl.cell(r,c); cell.text=str(val); cell.vertical_anchor=MSO_ANCHOR.MIDDLE; cell.fill.solid(); cell.fill.fore_color.rgb=BLUE if r==0 else (LIGHT if r%2 else WHITE)
            for p in cell.text_frame.paragraphs:
                p.alignment=PP_ALIGN.CENTER; p.font.name=FONT; p.font.size=Pt(size); p.font.bold=(r==0); p.font.color.rgb=WHITE if r==0 else NAVY

def flow(slide, labels,x,y,bw,gap,color=BLUE):
    for i,lab in enumerate(labels):
        bx=x+i*(bw+gap); s=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(bx),Inches(y),Inches(bw),Inches(0.65)); s.fill.solid(); s.fill.fore_color.rgb=LIGHT; s.line.color.rgb=color
        p=s.text_frame.paragraphs[0]; p.text=lab; p.font.name=FONT; p.font.size=Pt(12); p.font.bold=True; p.font.color.rgb=NAVY; p.alignment=PP_ALIGN.CENTER
        if i<len(labels)-1:
            a=slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(bx+bw+0.03),Inches(y+0.18),Inches(gap-0.06),Inches(0.28)); a.fill.solid(); a.fill.fore_color.rgb=color; a.line.fill.background()

slides=[]
def new(t,sub=None):
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s,t,sub); slides.append(s); return s

# 1
s=prs.slides.add_slide(prs.slide_layouts[6]); slides.append(s)
sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height); sh.fill.solid(); sh.fill.fore_color.rgb=NAVY; sh.line.fill.background()
bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,Inches(0.18),prs.slide_height); bar.fill.solid(); bar.fill.fore_color.rgb=CYAN; bar.line.fill.background()
b=s.shapes.add_textbox(Inches(0.9),Inches(1.4),Inches(11.5),Inches(1.2)); p=b.text_frame.paragraphs[0]; p.text='STF³-Detect'; p.font.name='Arial'; p.font.size=Pt(54); p.font.bold=True; p.font.color.rgb=WHITE
b=s.shapes.add_textbox(Inches(0.95),Inches(2.55),Inches(11.6),Inches(0.8)); p=b.text_frame.paragraphs[0]; p.text='基于时空-频域三分支特征融合的 AI 生成视频检测系统'; p.font.name=FONT; p.font.size=Pt(28); p.font.color.rgb=WHITE
b=s.shapes.add_textbox(Inches(0.98),Inches(4.15),Inches(10.5),Inches(0.6)); p=b.text_frame.paragraphs[0]; p.text='内容安全 Final Project 开题答辩'; p.font.name=FONT; p.font.size=Pt(18); p.font.color.rgb=RGBColor(210,230,255)
b=s.shapes.add_textbox(Inches(0.98),Inches(5.15),Inches(10.5),Inches(0.5)); p=b.text_frame.paragraphs[0]; p.text='数据集：GenVideo-Val  |  实现版本：STF³-Detect-Lite'; p.font.name=FONT; p.font.size=Pt(15); p.font.color.rgb=RGBColor(210,230,255)

s=new('01 选题背景：AI 生成视频正在快速逼近真实视频')
bullets(s,['文生视频 / 图生视频模型快速发展：Sora、Kling、Runway、Pika、OpenSora 等','生成视频已经具备较强真实感，单帧肉眼判断越来越困难','传统 DeepFake 多关注人脸篡改，而 AI 生成视频往往是整段内容合成'],0.8,1.35,6.0,2.4,18)
card(s,7.1,1.35,2.4,1.5,'Spatial','局部纹理、边缘、上采样伪影',BLUE); card(s,9.8,1.35,2.4,1.5,'Temporal','帧间闪烁、形变、运动不连续',CYAN); card(s,8.45,3.2,2.4,1.5,'Frequency','高频统计、频谱分布异常',ORANGE)
bullets(s,['核心判断：视频真假检测需要从“单帧 + 帧间 + 频谱”共同观察。'],0.9,5.45,11.5,0.6,17,BLUE)

s=new('02 内容安全意义：AI 生成视频带来的风险')
card(s,0.8,1.35,2.8,1.35,'虚假信息传播','生成虚假新闻现场、灾害现场、社会事件视频',RED); card(s,3.9,1.35,2.8,1.35,'诈骗与冒充','生成高管、亲友或公众人物视频用于欺诈',ORANGE); card(s,7.0,1.35,2.8,1.35,'证据可信度下降','伪造看似真实的视频证据，影响司法和媒体判断',BLUE); card(s,10.1,1.35,2.4,1.35,'平台治理压力','AIGC 内容规模增长，人工审核成本高',CYAN)
bullets(s,['目标：构建一个可运行、可解释、可展示的 AI 生成视频检测系统','场景：辅助内容审核、风险提示、课程内容安全实验展示'],1.0,3.8,10.8,1.4,20)

s=new('03 研究挑战')
card(s,0.8,1.3,3.6,1.3,'挑战 1：单帧真实感强','单帧分类容易受到语义和场景影响',BLUE); card(s,4.85,1.3,3.6,1.3,'挑战 2：时序伪影隐蔽','结构闪烁、动作不连续可能只出现在局部帧间变化中',CYAN); card(s,8.9,1.3,3.6,1.3,'挑战 3：跨生成器泛化','不同生成器伪影不同，模型不能只记住某一类生成器',ORANGE)
bullets(s,['本项目的切入点：融合空间纹理、时序变化、频域统计三类互补线索','项目边界：开题阶段只使用 GenVideo-Val，先验证完整 pipeline 可行性'],1.0,4.1,11.2,1.4,19)

s=new('04 相关工作与参考基础')
table(s,[['方向','代表工作','启发'],['AI 图像检测','UniversalFakeDetect / NPR / FreqNet','空间纹理与频域伪影可用于检测'],['DeepFake 视频检测','Face X-ray / SBI / AVFF','视频检测需考虑时序，但人脸依赖较强'],['AI 生成视频检测','DeMamba / D3 / GenVidBench','关注 AI 视频生成器伪影和跨生成器泛化']],0.75,1.3,11.85,2.2,11)
bullets(s,['核心参考代码：DeMamba、NPR-DeepfakeDetection、FreqNet、D3','本项目不完整复现大模型，而是实现轻量可运行版本 STF³-Detect-Lite'],0.9,4.2,11.5,1.1,18)

s=new('05 数据集：GenVideo-Val 单数据集方案')
table(s,[['类别','视频数'],['Real / real_MSRVTT','10000'],['Fake total','8302'],['总计','18302']],0.9,1.35,4.0,2.0,13)
bullets(s,['文件大小约 13.93GB，适合本机 RTX 4060 环境完成实验','包含真实视频和多种 AI 生成器视频，可支持二分类和跨生成器测试','已完成目录整理、metadata.csv、random split 与 OOD split 生成'],5.4,1.35,6.7,2.2,17)
card(s,0.9,4.35,3.4,1.0,'Random Split','70% train / 15% val / 15% test',BLUE); card(s,4.8,4.35,3.4,1.0,'OOD Split','MorphStudio / Show_1 / Sora / WildScrape 作为未见生成器',ORANGE); card(s,8.7,4.35,3.4,1.0,'Mini Split','开题可行性验证：320/80/160',GREEN)

s=new('06 方法总览：STF³-Detect-Lite')
flow(s,['输入视频','采样 16/8/4 帧','三分支特征','特征融合','真假分类'],0.8,1.25,2.0,0.45,BLUE)
card(s,0.95,2.75,3.5,1.4,'空间分支 Spatial','ResNet18 提取单帧纹理与局部结构特征；对时间维做平均池化。',BLUE); card(s,4.95,2.75,3.5,1.4,'时序分支 Temporal','计算相邻帧差分，捕捉闪烁、形变和运动不连续。',CYAN); card(s,8.95,2.75,3.5,1.4,'频域分支 Frequency','FFT log-magnitude spectrum + 小 CNN，捕捉频谱统计异常。',ORANGE)
bullets(s,['融合方式：Concat(Fs, Ft, Ff) → MLP → Real / AI-generated'],1.0,5.55,11,0.5,18,BLUE)

s=new('07 三分支设计细节')
table(s,[['分支','输入','实现','输出'],['Spatial','RGB frames','ResNet18 + temporal mean','F_s'],['Temporal','frame differences','Small CNN + mean','F_t'],['Frequency','FFT spectrum','Small CNN + mean','F_f'],['Fusion','F_s, F_t, F_f','Concat + MLP','Real/Fake logits']],0.8,1.3,11.8,3.0,12)
bullets(s,['轻量化原因：RTX 4060 8GB 显存，优先保证可训练、可复现、可展示','RAFT / Mamba / CLIP 可作为后续扩展项，不作为开题阶段必须实现'],0.95,4.8,11.5,1.0,17)

s=new('08 实验设计')
card(s,0.8,1.35,3.7,1.35,'Baseline 对比','Spatial、Frequency、Temporal、Spatial+Frequency、STF³',BLUE); card(s,4.85,1.35,3.7,1.35,'消融实验','验证空间、时序、频域各分支及融合策略贡献',CYAN); card(s,8.9,1.35,3.7,1.35,'跨生成器测试','使用 OOD split 评估遇到未见生成器时的泛化能力',ORANGE)
table(s,[['指标','含义'],['ACC','整体分类准确率'],['AUC','阈值无关的二分类检测能力'],['F1 / Precision / Recall','误报与漏报分析'],['Confusion Matrix','错误类型可视化']],1.0,3.45,11.2,2.15,12)

s=new('09 初步可行性验证：Mini Experiment','注：小样本实验仅用于开题阶段验证 pipeline，不代表最终完整实验性能')
table(s,[['项目','设置'],['训练集','320 videos，real/fake=160/160'],['验证集','80 videos，real/fake=40/40'],['测试集','160 videos，real/fake=80/80'],['采样/分辨率','4 frames, 112×112'],['训练轮数','2 epochs']],0.8,1.35,5.0,3.15,12)
pic(s,ASSETS/'mini_stf3_frequency_examples.png',6.25,1.35,w=5.7)
bullets(s,['验证内容：数据读取、视频采样、训练、评估、可视化、Demo 均已跑通。'],0.95,5.35,11.4,0.45,17,BLUE)

s=new('10 Mini Experiment 结果：三分支融合更有潜力')
table(s,[['方法','ACC','AUC','F1','Precision','Recall'],['Spatial','0.5563','0.6641','0.6537','0.5360','0.8375'],['Frequency','0.5500','0.6547','0.6400','0.5333','0.8000'],['STF³-Lite','0.7625','0.8764','0.7164','0.8889','0.6000']],0.7,1.25,6.0,2.3,11)
pic(s,ASSETS/'mini_results_bar.png',7.05,1.2,w=5.4)
bullets(s,['小样本下，STF³-Detect-Lite 的 AUC 达到 0.8764，高于单分支 baseline','说明三类特征具有互补性，后续完整实验值得继续推进'],0.9,4.55,11.5,1.0,17)

s=new('11 可视化结果：ROC 与混淆矩阵')
pic(s,ASSETS/'mini_stf3_roc_curve.png',0.9,1.25,w=5.4); pic(s,ASSETS/'mini_stf3_confusion_matrix.png',7.0,1.25,w=4.6)
bullets(s,['ROC 曲线展示模型在不同阈值下的检测能力','混淆矩阵可用于分析误报与漏报：开题阶段重点是验证流程和方向'],0.95,5.6,11.6,0.7,16)

s=new('12 Demo 原型设计')
flow(s,['上传视频','均匀采样','模型推理','输出概率','展示采样帧/频谱'],0.75,1.25,2.05,0.38,CYAN)
card(s,0.9,2.55,3.5,1.35,'输入','mp4 / mov 等常见视频格式',BLUE); card(s,4.9,2.55,3.5,1.35,'输出','Real 概率、AI-generated 概率、预测标签',GREEN); card(s,8.9,2.55,3.5,1.35,'展示','采样帧预览 + FFT 频谱图',ORANGE)
bullets(s,['已实现 Gradio Demo：python -m src.demo --checkpoint runs\\mini_stf3\\best.pt','答辩策略：优先准备录屏和截图，现场演示作为加分项'],1.0,5.05,11.4,1.0,17)

s=new('13 当前实现进度')
table(s,[['模块','状态'],['环境搭建 / GPU 验证','已完成'],['GenVideo-Val 解压整理','已完成'],['metadata 与 split CSV','已完成'],['Dataset + 视频采样','已完成'],['三分支模型','已完成 Lite 版'],['训练/评估/可视化脚本','已完成'],['Gradio Demo','已完成原型']],0.8,1.25,6.4,4.2,12)
bullets(s,['本机配置：RTX 4060 Laptop GPU，约 8GB 显存','核心环境：PyTorch 2.11.0 + CUDA 12.8','项目路径：final_project/src + scripts + docs/ppt_assets'],7.65,1.55,4.9,2.3,16)

s=new('14 后续计划')
table(s,[['阶段','任务','产出'],['1','完整 baseline 训练','Spatial/Frequency/Temporal 指标'],['2','三分支模型优化','STF³ 完整实验结果'],['3','OOD 跨生成器测试','泛化性分析表格'],['4','可视化与错误案例','ROC/混淆矩阵/频谱图'],['5','Demo 与报告整理','答辩 PPT、报告、录屏']],0.8,1.35,11.8,3.3,12)
bullets(s,['优先级：先保证完整实验闭环，再提升模型复杂度和指标。'],1.0,5.3,11.0,0.5,18,BLUE)

s=new('15 风险评估与应对')
table(s,[['风险','影响','应对'],['显存不足','训练 OOM','降低 batch size / 分辨率 / 帧数'],['单数据集泛化有限','结论说服力弱','增加 OOD split 分析'],['时序大模型复现难','进度风险','使用帧间差分 Lite 分支'],['Demo 现场不稳定','展示风险','准备录屏和固定样例']],0.8,1.35,11.8,3.0,12)
bullets(s,['项目策略：以可完成、可解释、可展示为主；复杂模型作为扩展方向。'],1.0,5.1,11.0,0.6,18,BLUE)

s=new('16 预期成果')
card(s,0.9,1.3,3.5,1.2,'代码','数据处理、Dataset、三分支模型、训练评估脚本',BLUE); card(s,4.9,1.3,3.5,1.2,'实验','Random split、OOD split、baseline、消融实验',CYAN); card(s,8.9,1.3,3.5,1.2,'展示','图表、Demo、报告、答辩 PPT、录屏',ORANGE)
bullets(s,['最终目标：形成一个完整的 AI 生成视频检测课程项目闭环','核心亮点：空间 + 时序 + 频域三类线索融合，兼顾可行性与可解释性'],1.0,3.75,11.2,1.2,19)

s=new('17 总结')
bullets(s,['选题面向 AI 生成视频带来的内容安全风险，具有现实意义','方法采用 STF³ 三分支融合：Spatial、Temporal、Frequency','已完成 GenVideo-Val 数据准备、模型代码、训练评估与 Demo 原型','小样本实验显示三分支融合优于单分支 baseline，验证了后续研究可行性'],1.0,1.5,11.0,3.0,21)
card(s,3.0,5.35,7.3,0.85,'下一步','扩大到完整 GenVideo-Val 实验，完成消融、OOD 泛化和最终 Demo 展示。',GREEN)

s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s); slides.append(s)
b=s.shapes.add_textbox(Inches(0),Inches(2.25),prs.slide_width,Inches(1.0)); p=b.text_frame.paragraphs[0]; p.text='Q & A'; p.font.name='Arial'; p.font.size=Pt(54); p.font.bold=True; p.font.color.rgb=NAVY; p.alignment=PP_ALIGN.CENTER
b=s.shapes.add_textbox(Inches(0),Inches(3.45),prs.slide_width,Inches(0.5)); p=b.text_frame.paragraphs[0]; p.text='谢谢老师和同学们'; p.font.name=FONT; p.font.size=Pt(24); p.font.color.rgb=GRAY; p.alignment=PP_ALIGN.CENTER

for i,s in enumerate(slides,1):
    if i!=1: page(s,i)
prs.save(OUT)
print('[write]',OUT,OUT.stat().st_size)
