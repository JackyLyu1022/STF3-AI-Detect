"""Inspect teacher's pptx template (Chinese-safe)."""
from pptx import Presentation
from pathlib import Path

src = "C:/Users/jacky/AppData/Local/Temp/teacher_template.pptx"
out = Path("docs/teacher_template_dump.txt")
p = Presentation(src)

lines = []
lines.append(f"width(in)={p.slide_width/914400:.2f} height(in)={p.slide_height/914400:.2f}")
lines.append(f"total slides: {len(p.slides)}")
for i, s in enumerate(p.slides, 1):
    lines.append("")
    lines.append(f"=== Slide {i} (layout={s.slide_layout.name}) ===")
    for sh in s.shapes:
        try:
            x = sh.left / 914400
            y = sh.top / 914400
            w = sh.width / 914400
            h = sh.height / 914400
        except Exception:
            x = y = w = h = 0
        kind = str(sh.shape_type)
        text = ""
        if sh.has_text_frame:
            text = sh.text_frame.text.replace("\n", " // ")[:300]
        lines.append(f"  [{kind}] x={x:.2f} y={y:.2f} w={w:.2f} h={h:.2f}")
        if text:
            lines.append(f"    text={text}")
out.write_text("\n".join(lines), encoding="utf-8")
print(f"wrote {out}")
